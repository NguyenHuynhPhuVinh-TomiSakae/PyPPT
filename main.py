import os
import docx
import google.generativeai as genai
from pptx import Presentation
from pptx.util import Inches
import random

# Cấu hình API Key cho Google Generative AI
genai.configure(api_key="AIzaSyB_eNpMTroPTupXzl_oey08M0d-luxJ3OE")

# Đọc nội dung từ tệp Word
def read_docx(file_path):
    doc = docx.Document(file_path)
    content = {}
    current_title = None
    
    for para in doc.paragraphs:
        if para.style.name == 'Heading 1':
            current_title = para.text
            content[current_title] = []
        elif current_title:
            content[current_title].append(para.text)
    
    return content

# Tóm tắt nội dung bằng AI
def summarize_text(text):
    model = genai.GenerativeModel(
        model_name="gemini-1.5-flash",
        generation_config={
            "temperature": 1,
            "top_p": 0.95,
            "top_k": 64,
            "max_output_tokens": 8192,
            "response_mime_type": "text/plain",
        },
    )
    
    chat_session = model.start_chat(
  history=[
    {
      "role": "user",
      "parts": [
        "Hãy tóm tắt toàn bộ nội dung và tôi đưa vào, các nội dung mà tôi đưa vào sẽ liên quan tới môn điện toán đám mây",
      ],
    },
    {
      "role": "model",
      "parts": [
        "Xin lỗi, bạn cần cung cấp thêm thông tin để tôi có thể tóm tắt nội dung của bạn. \n\nBạn muốn tôi tóm tắt nội dung gì về điện toán đám mây? Ví dụ:\n\n* **Bài viết/ tài liệu cụ thể:** Vui lòng chia sẻ bài viết hoặc tài liệu bạn muốn tóm tắt.\n* **Chủ đề cụ thể:** Bạn muốn tóm tắt nội dung về một chủ đề cụ thể trong điện toán đám mây (ví dụ: các loại dịch vụ đám mây, kiến trúc đám mây, bảo mật đám mây,...)?\n* **Câu hỏi cụ thể:** Bạn muốn tôi tóm tắt câu trả lời cho một câu hỏi cụ thể về điện toán đám mây?\n\nCung cấp thêm thông tin chi tiết để tôi có thể hỗ trợ bạn tốt hơn.\n",
      ],
    },
    {
      "role": "user",
      "parts": [
        "Chương trước mô tả bức tường năng lượng buộc ngành điện toán phải chuyển sang sử dụng nhiều bộ xử lý và nhiều máy tính. Chương này cũng thảo luận về cách thu thập máy chủ vào một cơ sở tập trung có thể giảm cả chi phí thiết bị và chi phí định kỳ của nhân viên CNTT.\nChương này giải thích cách phần mềm cho phép sử dụng một bộ máy chủ theo cách hoàn toàn mới. Chương này mô tả khái niệm chính của điện toán đàn hồi và giải thích cách ảo hóa cho phép chủ sở hữu đám mây mở rộng quy mô dịch vụ.\nChương này cũng mô tả các mô hình kinh doanh phát sinh trong ngành điện toán đám mây.\nChương tiếp theo mô tả cách các công ty điện toán đám mây công cộng phát sinh sử dụng điện toán đàn hồi. Các chương sau giải thích cơ sở hạ tầng và công nghệ cơ bản cho phép điện toán đàn hồi cũng như các hệ thống phần mềm được sử dụng để quản lý các triển khai trong đám mây.\n",
      ],
    },
    {
      "role": "model",
      "parts": [
        "Chương này thảo luận về cách ngành điện toán đã phải thay đổi để đối phó với giới hạn của bức tường năng lượng. Thay vì sử dụng nhiều bộ xử lý và máy tính, ngành này đã chuyển sang tập trung các máy chủ vào một cơ sở tập trung. \n\nChương này giải thích cách phần mềm cho phép tận dụng các máy chủ theo cách hoàn toàn mới, thông qua khái niệm \"điện toán đàn hồi\". Ảo hóa đóng vai trò quan trọng trong việc cho phép chủ sở hữu đám mây mở rộng quy mô dịch vụ của họ. \n\nChương này cũng mô tả các mô hình kinh doanh mới phát sinh từ ngành điện toán đám mây. \n\nChương tiếp theo sẽ đi sâu vào cách các công ty điện toán đám mây công cộng tận dụng điện toán đàn hồi, đồng thời khám phá cơ sở hạ tầng và công nghệ hỗ trợ nó, cũng như các hệ thống phần mềm quản lý các triển khai trong đám mây. \n",
      ],
    },
    {
      "role": "user",
      "parts": [
        "Hãy rút gọn trong 1 dòng",
      ],
    },
    {
      "role": "model",
      "parts": [
        "Chương này giải thích cách ngành điện toán đã chuyển sang điện toán đàn hồi, sử dụng ảo hóa để mở rộng quy mô dịch vụ và tạo ra các mô hình kinh doanh mới. \n",
      ],
    },
    {
      "role": "user",
      "parts": [
        "rút gọn 1 dòng nhưng không quá ngắn để phù hợp cho power point",
      ],
    },
    {
      "role": "model",
      "parts": [
        "Chương này giới thiệu khái niệm điện toán đàn hồi, cho thấy cách nó giải quyết giới hạn của bức tường năng lượng và tạo ra các mô hình kinh doanh mới trong ngành điện toán. \n",
      ],
    },
    {
      "role": "user",
      "parts": [
        "dài hơn",
      ],
    },
    {
      "role": "model",
      "parts": [
        "Chương này trình bày cách ngành điện toán đã đối mặt và vượt qua giới hạn của \"bức tường năng lượng\" bằng cách chuyển sang điện toán đàn hồi. Nó thảo luận về vai trò của ảo hóa trong việc cho phép chủ sở hữu đám mây mở rộng quy mô dịch vụ một cách hiệu quả, đồng thời giới thiệu các mô hình kinh doanh mới phát triển từ sự bùng nổ của điện toán đám mây. \n",
      ],
    },
    {
      "role": "user",
      "parts": [
        "dài hơn nữa",
      ],
    },
    {
      "role": "model",
      "parts": [
        "Chương này khai thác vấn đề \"bức tường năng lượng\" đã từng là rào cản cho ngành điện toán, dẫn đến việc chuyển đổi sang sử dụng nhiều bộ xử lý và máy tính. Nó giải thích cách tập trung máy chủ vào một cơ sở tập trung, kết hợp với phần mềm ảo hóa, đã tạo ra khái niệm \"điện toán đàn hồi\". Chương này đi sâu vào cách điện toán đàn hồi cho phép chủ sở hữu đám mây mở rộng quy mô dịch vụ một cách linh hoạt và hiệu quả hơn, đồng thời khám phá các mô hình kinh doanh mới đã xuất hiện trong ngành điện toán đám mây. \n",
      ],
    },
  ]
)
    
    prompts = [
        "Tóm tắt nội dung chính của chương này trong một câu dài và chi tiết: ",
        "Trình bày các điểm quan trọng nhất của chương này trong một đoạn ngắn gọn: ",
        "Hãy tóm lược chương này thành 3-4 ý chính: ",
        "Tổng hợp nội dung của chương này thành một đoạn súc tích nhưng đầy đủ thông tin: ",
        "Trình bày các khái niệm chính được đề cập trong chương này: ",
        "Hãy tóm tắt chương này dưới dạng một danh sách các điểm chính: ",
        "Trình bày nội dung chính của chương này dưới dạng một sơ đồ tư duy ngắn gọn: ",
        "Tóm tắt chương này bằng cách nhấn mạnh vào các ứng dụng và tác động của điện toán đám mây: ",
        "Tổng hợp nội dung chính của chương này dưới dạng một tin tức ngắn về sự phát triển của ngành điện toán: ",
        "Hãy tóm tắt chương này bằng cách nhấn mạnh vào lợi ích của điện toán đám mây đối với doanh nghiệp: ",
        "Tóm tắt chương này dưới dạng một bài báo tin tức ngắn gọn: ",
        "Tạo một bảng so sánh ngắn gọn giữa các khái niệm chính được đề cập trong chương này: ",
        "Viết một đoạn văn ngắn giải thích nội dung chương này cho một người không có kiến thức về công nghệ: ",
    ]
    chosen_prompt = random.choice(prompts)
    response = chat_session.send_message(chosen_prompt + text)
    return response.text

# Tạo bài thuyết trình PowerPoint từ nội dung tóm tắt
def create_presentation(content, output_path):
    # Tìm id mới cho file output
    new_id = 1
    while os.path.exists(f'output/output{new_id}.pptx'):
        new_id += 1
    
    output_path = f'output/output{new_id}.pptx'
    
    prs = Presentation()
    prs.save(output_path)
    
    for title, paragraphs in content.items():
        # Kết hợp các đoạn văn thành một văn bản duy nhất
        full_text = '\n'.join(paragraphs)
        print("Đang tạo Slide: " + title)
        # Tóm tắt nội dung
        summary = summarize_text(full_text)
        print("Tạo Slide: " + title +" thành công")
        # Thêm slide tiêu đề
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Slide layout 1: Title and Content
        title_shape = slide.shapes.title
        body_shape = slide.placeholders[1]
        
        title_shape.text = title
        body_shape.text = summary
    
    prs.save(output_path)
    return output_path

# Đọc nội dung từ file data.docx
docx_path = 'data.docx'
content = read_docx(docx_path)

# Tạo bài thuyết trình PowerPoint
pptx_path = create_presentation(content, 'output.pptx')

print("Bài thuyết trình đã được lưu thành công.")
print(f"Hãy tải về theo đường dẫn sau:")
print(f"https://da22ttc-tvu.github.io/pptx?name={os.path.basename(pptx_path)}")